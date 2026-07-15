"""her-desktop office sidecar: pptx/docx/xlsx tooling.

Usage: office_tool.py <tool> <payload.json>
Reads a JSON payload file, prints a JSON object to stdout:
  {"ok": true, ...tool fields...} or {"error": "message"}.
Tools: ppt-generate, docx-generate, docx-read, xlsx-read, xlsx-write.
"""
import json
import sys


def fail(message):
    print(json.dumps({"error": str(message)[:500]}, ensure_ascii=False))
    sys.exit(0)


def ppt_generate(payload):
    from pptx import Presentation
    from pptx.util import Pt

    slides = payload.get("slides") or []
    if not slides:
        fail("slides 不能为空：需要 [{title, bullets[], notes?}, …]")
    prs = Presentation()
    deck_title = payload.get("title") or ""
    if deck_title:
        layout = prs.slide_layouts[0]  # title slide
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = deck_title
        subtitle = payload.get("subtitle") or ""
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
    for spec in slides:
        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(spec.get("title") or "")
        bullets = spec.get("bullets") or []
        body = slide.placeholders[1].text_frame
        for index, bullet in enumerate(bullets):
            text = str(bullet)
            level = 0
            # 两个空格缩进 = 下一级
            while text.startswith("  "):
                level += 1
                text = text[2:]
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = text.strip()
            paragraph.level = min(level, 4)
            paragraph.font.size = Pt(20 if level == 0 else 16)
        notes = spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = str(notes)
    prs.save(payload["output_path"])
    return {"ok": True, "path": payload["output_path"], "slide_count": len(prs.slides._sldIdLst)}


def docx_generate(payload):
    import re

    from docx import Document

    markdown = payload.get("markdown") or ""
    if not markdown.strip():
        fail("markdown 不能为空")
    document = Document()
    title = payload.get("title") or ""
    if title:
        document.add_heading(title, level=0)
    bold = re.compile(r"\*\*(.+?)\*\*")

    def add_runs(paragraph, text):
        pos = 0
        for match in bold.finditer(text):
            if match.start() > pos:
                paragraph.add_run(text[pos:match.start()])
            paragraph.add_run(match.group(1)).bold = True
            pos = match.end()
        if pos < len(text):
            paragraph.add_run(text[pos:])

    for raw in markdown.splitlines():
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            level = min(len(stripped) - len(stripped.lstrip("#")), 4)
            document.add_heading(stripped.lstrip("#").strip(), level=level)
        elif stripped.startswith(("- ", "* ")):
            add_runs(document.add_paragraph(style="List Bullet"), stripped[2:])
        elif re.match(r"^\d+[.)] ", stripped):
            add_runs(document.add_paragraph(style="List Number"), re.sub(r"^\d+[.)] ", "", stripped))
        else:
            add_runs(document.add_paragraph(), stripped)
    document.save(payload["output_path"])
    return {"ok": True, "path": payload["output_path"]}


def docx_read(payload):
    from docx import Document

    document = Document(payload["path"])
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    text = "\n".join(parts)
    limit = int(payload.get("max_chars") or 20000)
    truncated = len(text) > limit
    return {"ok": True, "text": text[:limit], "truncated": truncated}


def xlsx_read(payload):
    from openpyxl import load_workbook

    workbook = load_workbook(payload["path"], read_only=True, data_only=True)
    sheet_name = payload.get("sheet")
    sheet = workbook[sheet_name] if sheet_name else workbook.active
    max_rows = int(payload.get("max_rows") or 60)
    lines = []
    total = 0
    for row in sheet.iter_rows(values_only=True):
        total += 1
        if total <= max_rows:
            lines.append(" | ".join("" if v is None else str(v) for v in row))
    return {
        "ok": True,
        "sheet": sheet.title,
        "sheets": workbook.sheetnames,
        "total_rows": total,
        "text": "\n".join(lines),
        "truncated": total > max_rows,
    }


def xlsx_write(payload):
    from openpyxl import Workbook

    rows = payload.get("rows")
    if not isinstance(rows, list) or not rows:
        fail("rows 不能为空：需要二维数组 [[单元格, …], …]")
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = str(payload.get("sheet_name") or "Sheet1")[:31]
    for row in rows:
        if not isinstance(row, list):
            fail("rows 必须是二维数组，每行是一个数组")
        sheet.append(["" if v is None else v for v in row])
    workbook.save(payload["output_path"])
    return {"ok": True, "path": payload["output_path"], "row_count": len(rows)}


TOOLS = {
    "ppt-generate": ppt_generate,
    "docx-generate": docx_generate,
    "docx-read": docx_read,
    "xlsx-read": xlsx_read,
    "xlsx-write": xlsx_write,
}


def main():
    if len(sys.argv) != 3 or sys.argv[1] not in TOOLS:
        fail(f"usage: office_tool.py <{'|'.join(TOOLS)}> <payload.json>")
    try:
        payload = json.load(open(sys.argv[2], encoding="utf-8"))
    except Exception as exc:  # noqa: BLE001
        fail(f"payload 解析失败: {exc}")
    try:
        print(json.dumps(TOOLS[sys.argv[1]](payload), ensure_ascii=False))
    except Exception as exc:  # noqa: BLE001
        fail(exc)


if __name__ == "__main__":
    main()
